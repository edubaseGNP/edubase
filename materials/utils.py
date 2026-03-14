"""
Utility functions for materials: image compression and OCR text extraction.
"""

import io
import logging
from typing import NamedTuple

logger = logging.getLogger(__name__)


class OCRResult(NamedTuple):
    text: str
    model_name: str = ''
    tokens_used: int = 0
    backend_used: str = ''
    file_type: str = 'unknown'


# ---------------------------------------------------------------------------
# Image compression
# ---------------------------------------------------------------------------

def compress_image_file(file_field, max_width: int = 1920, quality: int = 85) -> None:
    """
    Resize and compress an uploaded image in-place.
    - If width > max_width: downscale proportionally.
    - JPEG for non-transparent images, PNG otherwise.
    """
    from PIL import Image

    try:
        file_field.seek(0)
        img = Image.open(file_field)
        has_alpha = img.mode in ('RGBA', 'LA', 'P')

        if img.width > max_width:
            ratio = max_width / img.width
            img = img.resize((max_width, int(img.height * ratio)), Image.LANCZOS)
            logger.debug('Image resized to %dx%d', img.width, img.height)

        output = io.BytesIO()
        if has_alpha:
            img.save(output, format='PNG', optimize=True)
            new_name = _replace_extension(file_field.name, '.png')
        else:
            img.convert('RGB').save(output, format='JPEG', quality=quality, optimize=True)
            new_name = _replace_extension(file_field.name, '.jpg')

        output.seek(0)
        file_field.save(new_name, output, save=False)
        logger.info('Image compressed: %s', new_name)

    except Exception:
        logger.exception('Image compression failed for %s', file_field.name)


def _replace_extension(filename: str, new_ext: str) -> str:
    import os
    base, _ = os.path.splitext(filename)
    return base + new_ext


# ---------------------------------------------------------------------------
# OCR helpers
# ---------------------------------------------------------------------------

def _preprocess_for_ocr(img):
    """
    Prepare a PIL image for Tesseract:
      1. Greyscale
      2. Auto-contrast
      3. Upscale to at least 1800 px wide (≈300 DPI for A4 scans)
      4. Sharpen
    Returns a new PIL image.
    """
    from PIL import Image, ImageFilter, ImageOps

    img = img.convert('L')                          # greyscale
    img = ImageOps.autocontrast(img, cutoff=2)      # stretch histogram, clip 2 % outliers

    if img.width < 1800:                            # upscale small/low-DPI scans
        scale = 1800 / img.width
        img = img.resize(
            (int(img.width * scale), int(img.height * scale)),
            Image.LANCZOS,
        )
        logger.debug('OCR preprocess: upscaled to %dx%d', img.width, img.height)

    img = img.filter(ImageFilter.SHARPEN)
    return img


def _looks_like_real_text(text: str, min_chars: int = 30) -> bool:
    """
    Sanity-check pdfminer output.
    Returns False if text is too short or mostly non-printable (garbled encoding).
    """
    if not text or len(text.strip()) < min_chars:
        return False
    printable = sum(1 for c in text if c.isprintable() and not c.isspace())
    ratio = printable / max(len(text), 1)
    return ratio > 0.65


# ---------------------------------------------------------------------------
# OCR / text extraction — public API
# ---------------------------------------------------------------------------

def _get_ai_cfg():
    """
    Return the SiteConfig instance (DB-stored AI settings).
    Falls back gracefully if DB is unavailable (e.g. during migrations).
    """
    try:
        from core.models import SiteConfig
        return SiteConfig.get()
    except Exception:
        return None


def _log_ai_call(
    backend: str,
    material_id,
    success: bool,
    chars: int,
    duration_ms: int,
    error: str = '',
    model_name: str = '',
    tokens_used: int = 0,
    file_type: str = 'unknown',
    attempt: int = 1,
    trigger: str = 'upload',
) -> None:
    """Persist one AICallLog record; silently skips on any error."""
    try:
        from core.models import AICallLog
        AICallLog.objects.create(
            backend=backend,
            material_id=material_id,
            success=success,
            chars_extracted=chars,
            duration_ms=duration_ms,
            error_msg=error[:500],
            model_name=model_name,
            tokens_used=tokens_used,
            file_type=file_type,
            attempt=attempt,
            trigger=trigger,
        )
        if not success:
            _alert_on_high_error_rate(material_id)
    except Exception:
        logger.debug('_log_ai_call failed (non-critical)', exc_info=True)


def _alert_on_high_error_rate(material_id) -> None:
    """Create notification + send email if error rate in last hour exceeds 30%."""
    try:
        import datetime
        from django.utils import timezone
        from core.models import AICallLog, Notification

        hour_ago = timezone.now() - datetime.timedelta(hours=1)
        recent = AICallLog.objects.filter(timestamp__gte=hour_ago)
        total = recent.count()
        if total < 5:
            return
        failures = recent.filter(success=False).count()
        if failures / total < 0.30:
            return

        from django.contrib.auth import get_user_model
        User = get_user_model()
        msg = (
            f'AI/OCR: {failures}/{total} volání selhalo za poslední hodinu '
            f'({failures/total:.0%}). Zkontrolujte nastavení AI backendu.'
        )
        cutoff = timezone.now() - datetime.timedelta(minutes=30)
        for user in User.objects.filter(is_superuser=True):
            if not Notification.objects.filter(
                recipient=user,
                verb__startswith='AI/OCR:',
                created_at__gte=cutoff,
            ).exists():
                Notification.objects.create(
                    recipient=user,
                    verb=msg,
                    target_url='/admin/core/aicalllog/',
                )
        # Email alert if SMTP configured
        try:
            from core.models import SiteConfig
            cfg = SiteConfig.get()
            if cfg.email_notifications_enabled and cfg.smtp_host:
                from django.core.mail import send_mail
                admins = list(User.objects.filter(is_superuser=True).values_list('email', flat=True))
                if admins:
                    send_mail(
                        subject='[EduBase] Varování: AI/OCR vysoká chybovost',
                        message=msg + '\n\nZobrazte log: /admin/core/aicalllog/',
                        from_email=cfg.default_from_email or 'noreply@edubase.tech',
                        recipient_list=admins,
                        fail_silently=True,
                    )
        except Exception:
            pass
    except Exception:
        logger.debug('_alert_on_high_error_rate failed', exc_info=True)


def extract_text_from_image(filepath: str) -> OCRResult:
    cfg = _get_ai_cfg()
    backend = (cfg.ai_backend if cfg else None) or 'none'

    extractor = {
        'anthropic': _claude_vision_extract,
        'google':    _gemini_vision_extract,
        'ollama':    _ollama_vision_extract,
    }.get(backend)

    if extractor:
        result = extractor(filepath)
        if result.text:
            return result._replace(file_type='image')
        logger.warning('%s vision failed for %s – falling back to Tesseract', backend, filepath)

    result = _tesseract_extract(filepath)
    return result._replace(file_type='image')


def extract_text_from_pdf(filepath: str) -> OCRResult:
    text = _pdfminer_extract(filepath)
    if _looks_like_real_text(text):
        logger.debug('pdfminer OK for %s (%d chars)', filepath, len(text))
        return OCRResult(text=text, model_name='pdfminer', backend_used='pdfminer', file_type='pdf_text')

    logger.debug('pdfminer insufficient – OCR fallback for %s', filepath)

    cfg = _get_ai_cfg()
    if cfg and cfg.ai_backend == 'ollama':
        result = _ollama_pdf_extract(filepath)
        if result.text:
            return result._replace(file_type='pdf_ocr')
        logger.warning('Ollama PDF OCR failed for %s – falling back to Tesseract', filepath)

    result = _pdf_ocr_extract(filepath)
    return result._replace(file_type='pdf_ocr')


def extract_text_from_office(filepath: str) -> OCRResult:
    import os
    ext = os.path.splitext(filepath)[1].lower()
    dispatch = {
        '.docx': extract_text_from_docx,
        '.pptx': extract_text_from_pptx,
        '.xlsx': extract_text_from_xlsx,
        '.odt':  extract_text_from_odf,
        '.ods':  extract_text_from_odf,
        '.odp':  extract_text_from_odf,
    }
    fn = dispatch.get(ext)
    text = fn(filepath) if fn else ''
    return OCRResult(text=text, model_name='office', backend_used='office', file_type='office')


# ---------------------------------------------------------------------------
# Internal extraction backends
# ---------------------------------------------------------------------------

def _tesseract_extract(filepath: str) -> OCRResult:
    """Tesseract OCR with preprocessing and LSTM engine."""
    try:
        import pytesseract
        from PIL import Image

        cfg    = _get_ai_cfg()
        lang   = (cfg.ocr_lang if cfg else None) or 'ces+eng'
        config = '--oem 3 --psm 6'   # LSTM engine + uniform text block

        img  = Image.open(filepath)
        img  = _preprocess_for_ocr(img)
        text = pytesseract.image_to_string(img, lang=lang, config=config)
        logger.debug('Tesseract: %d chars from %s', len(text), filepath)
        return OCRResult(text=text.strip(), model_name='tesseract', tokens_used=0, backend_used='tesseract')
    except Exception:
        logger.exception('Tesseract failed for %s', filepath)
        return OCRResult(text='', backend_used='tesseract')


def _claude_vision_extract(filepath: str) -> OCRResult:
    """
    Use Claude Haiku vision to extract text from an image.
    Dramatically better for handwriting, photos of whiteboards, low-quality scans.
    Requires ANTHROPIC_API_KEY in environment.
    Cost: ~0.08 Kč / image (Haiku model).
    """
    import base64
    import pathlib

    try:
        import anthropic
    except ImportError:
        logger.error('anthropic package not installed – run: pip install anthropic')
        return OCRResult(text='', backend_used='anthropic')

    cfg     = _get_ai_cfg()
    api_key = (cfg.anthropic_api_key if cfg else '') or ''
    if not api_key:
        logger.error('Anthropic API key not set in SiteConfig')
        return OCRResult(text='', backend_used='anthropic')

    _MIME = {
        '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.png': 'image/png', '.webp': 'image/webp', '.gif': 'image/gif',
    }
    ext  = pathlib.Path(filepath).suffix.lower()
    mime = _MIME.get(ext, 'image/jpeg')

    try:
        data = base64.standard_b64encode(pathlib.Path(filepath).read_bytes()).decode()
        client = anthropic.Anthropic(api_key=api_key)
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=4096,
            messages=[{
                'role': 'user',
                'content': [
                    {
                        'type': 'image',
                        'source': {'type': 'base64', 'media_type': mime, 'data': data},
                    },
                    {
                        'type': 'text',
                        'text': (
                            'Extrahuj veškerý text z tohoto obrázku. '
                            'Zachovej strukturu, odrážky a nadpisy. '
                            'Vrať pouze čistý text bez úvodního komentáře.'
                        ),
                    },
                ],
            }],
        )
        text = msg.content[0].text.strip()
        tokens = (getattr(msg.usage, 'input_tokens', 0) or 0) + (getattr(msg.usage, 'output_tokens', 0) or 0)
        logger.info('Claude Vision: %d chars from %s', len(text), filepath)
        return OCRResult(text=text, model_name='claude-haiku-4-5-20251001', tokens_used=tokens, backend_used='anthropic')
    except Exception:
        logger.exception('Claude Vision failed for %s', filepath)
        return OCRResult(text='', backend_used='anthropic')


def _gemini_vision_extract(filepath: str) -> OCRResult:
    """
    Use Google Gemini Flash to extract text from an image.
    Free tier: 1 500 req/day. Paid: ~0.01 Kč/image (10× cheaper than Claude Haiku).
    Requires: pip install google-genai  +  GOOGLE_AI_API_KEY in env.
    """
    import base64
    import pathlib

    cfg     = _get_ai_cfg()
    api_key = (cfg.google_ai_api_key if cfg else '') or ''
    if not api_key:
        logger.error('GOOGLE_AI_API_KEY not set in SiteConfig')
        return OCRResult(text='', backend_used='google')

    try:
        from google import genai
        from google.genai import types
    except ImportError:
        logger.error('google-genai not installed – run: pip install google-genai')
        return OCRResult(text='', backend_used='google')

    _MIME = {
        '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.png': 'image/png',  '.webp': 'image/webp', '.gif': 'image/gif',
    }
    ext  = pathlib.Path(filepath).suffix.lower()
    mime = _MIME.get(ext, 'image/jpeg')

    model = (cfg.google_ai_model if cfg else '') or 'gemini-2.5-flash'

    try:
        data   = base64.standard_b64encode(pathlib.Path(filepath).read_bytes()).decode()
        client = genai.Client(api_key=api_key)
        resp   = client.models.generate_content(
            model=model,
            contents=[
                types.Part.from_bytes(data=base64.b64decode(data), mime_type=mime),
                'Extrahuj veškerý text z tohoto obrázku. '
                'Zachovej strukturu, odrážky a nadpisy. '
                'Vrať pouze čistý text bez úvodního komentáře.',
            ],
        )
        text = resp.text.strip()
        tokens = getattr(getattr(resp, 'usage_metadata', None), 'total_token_count', 0) or 0
        logger.info('Gemini Vision (%s): %d chars from %s', model, len(text), filepath)
        return OCRResult(text=text, model_name=model, tokens_used=tokens, backend_used='google')
    except Exception:
        logger.exception('Gemini Vision failed for %s', filepath)
        return OCRResult(text='', backend_used='google')


def _ollama_vision_extract(filepath: str) -> OCRResult:
    """
    Use a locally running Ollama server for vision OCR.
    Cost: 0 Kč – runs on your Proxmox server.
    Requires: Ollama running with llama3.2-vision (or similar vision model).
    docker-compose service 'ollama' or OLLAMA_BASE_URL env var.
    """
    import base64
    import pathlib
    import urllib.request
    import json

    cfg      = _get_ai_cfg()
    base_url = ((cfg.ollama_base_url if cfg else '') or 'http://ollama:11434').rstrip('/')
    model    = (cfg.ollama_vision_model if cfg else '') or 'llama3.2-vision'

    try:
        data = base64.standard_b64encode(pathlib.Path(filepath).read_bytes()).decode()
        payload = json.dumps({
            'model': model,
            'prompt': (
                'Extrahuj veškerý text z tohoto obrázku. '
                'Zachovej strukturu a odrážky. Vrať jen čistý text.'
            ),
            'images': [data],
            'stream': False,
        }).encode()
        req = urllib.request.Request(
            f'{base_url}/api/generate',
            data=payload,
            headers={'Content-Type': 'application/json'},
        )
        with urllib.request.urlopen(req, timeout=120) as resp:
            result = json.loads(resp.read())
        text = result.get('response', '').strip()
        logger.info('Ollama Vision (%s): %d chars from %s', model, len(text), filepath)
        return OCRResult(text=text, model_name=model, tokens_used=0, backend_used='ollama')
    except Exception:
        logger.exception('Ollama Vision failed for %s', filepath)
        return OCRResult(text='', backend_used='ollama')


def _ollama_pdf_extract(filepath: str) -> OCRResult:
    """
    Convert PDF pages to images and run Ollama Vision on each page.
    Used as OCR fallback when AI_BACKEND=ollama and pdfminer finds no real text.
    """
    import base64
    import json
    import urllib.request

    cfg      = _get_ai_cfg()
    base_url = ((cfg.ollama_base_url if cfg else '') or 'http://ollama:11434').rstrip('/')
    model    = (cfg.ollama_vision_model if cfg else '') or 'llama3.2-vision'
    dpi      = (cfg.ocr_pdf_dpi if cfg else None) or 300

    try:
        from pdf2image import convert_from_path
        pages = convert_from_path(filepath, dpi=dpi)
    except Exception:
        logger.exception('Ollama PDF: pdf2image failed for %s', filepath)
        return OCRResult(text='', backend_used='ollama')

    texts = []
    for i, page in enumerate(pages, 1):
        try:
            import io
            buf = io.BytesIO()
            page.save(buf, format='JPEG', quality=85)
            data = base64.standard_b64encode(buf.getvalue()).decode()
            payload = json.dumps({
                'model': model,
                'prompt': (
                    'Extrahuj veškerý text z této stránky dokumentu. '
                    'Zachovej strukturu, odrážky a nadpisy. Vrať jen čistý text.'
                ),
                'images': [data],
                'stream': False,
            }).encode()
            req = urllib.request.Request(
                f'{base_url}/api/generate',
                data=payload,
                headers={'Content-Type': 'application/json'},
            )
            with urllib.request.urlopen(req, timeout=180) as resp:
                result = json.loads(resp.read())
            texts.append(result.get('response', '').strip())
            logger.debug('Ollama PDF page %d/%d: %d chars', i, len(pages), len(texts[-1]))
        except Exception:
            logger.exception('Ollama PDF: page %d failed for %s', i, filepath)
            texts.append('')

    result = '\n\n'.join(t for t in texts if t).strip()
    logger.info('Ollama PDF (%s): %d chars, %d pages for %s', model, len(result), len(pages), filepath)
    return OCRResult(text=result, model_name=model, tokens_used=0, backend_used='ollama')


def _pdfminer_extract(filepath: str) -> str:
    try:
        from pdfminer.high_level import extract_text
        return extract_text(filepath) or ''
    except Exception:
        logger.exception('pdfminer failed for %s', filepath)
        return ''


def _pdf_ocr_extract(filepath: str) -> OCRResult:
    """Convert PDF pages to images at OCR_PDF_DPI and run Tesseract on each."""
    try:
        import pytesseract
        from pdf2image import convert_from_path

        cfg    = _get_ai_cfg()
        dpi    = (cfg.ocr_pdf_dpi if cfg else None) or 300
        lang   = (cfg.ocr_lang if cfg else None) or 'ces+eng'
        config = '--oem 3 --psm 6'

        pages = convert_from_path(filepath, dpi=dpi)
        texts = []
        for page in pages:
            page = _preprocess_for_ocr(page)
            texts.append(pytesseract.image_to_string(page, lang=lang, config=config))
        result = '\n\n'.join(texts).strip()
        logger.debug('PDF OCR (%d dpi): %d chars, %d pages', dpi, len(result), len(pages))
        return OCRResult(text=result, model_name='tesseract', tokens_used=0, backend_used='tesseract')
    except Exception:
        logger.exception('PDF OCR fallback failed for %s', filepath)
        return OCRResult(text='', backend_used='tesseract')


# ---------------------------------------------------------------------------
# Office document text extraction
# ---------------------------------------------------------------------------

def extract_text_from_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc   = Document(filepath)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_texts = [c.text.strip() for c in row.cells if c.text.strip()]
                if row_texts:
                    parts.append(' | '.join(row_texts))
        return '\n'.join(parts).strip()
    except Exception:
        logger.exception('docx extraction failed for %s', filepath)
        return ''


def extract_text_from_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
            if texts:
                parts.append(f'--- Snímek {i} ---\n' + '\n'.join(texts))
        return '\n\n'.join(parts).strip()
    except Exception:
        logger.exception('pptx extraction failed for %s', filepath)
        return ''


def extract_text_from_xlsx(filepath: str) -> str:
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None and str(v).strip()]
                if vals:
                    rows.append(' | '.join(vals))
            if rows:
                parts.append(f'=== {sheet.title} ===\n' + '\n'.join(rows))
        wb.close()
        return '\n\n'.join(parts).strip()
    except Exception:
        logger.exception('xlsx extraction failed for %s', filepath)
        return ''


def extract_text_from_odf(filepath: str) -> str:
    try:
        from odf.opendocument import load
        from odf import text as odftext
        from odf.teletype import extractText
        doc   = load(filepath)
        texts = [extractText(e).strip() for e in doc.body.getElementsByType(odftext.P)]
        return '\n'.join(t for t in texts if t).strip()
    except Exception:
        logger.exception('ODF extraction failed for %s', filepath)
        return ''
